#pylint: disable=invalid-name
"""
    Powerpoint module
    -----------------
    Interface for Arty to use the python-pptx library in order to
    generate Powerpoint presentations.
"""

import os
import re

import PIL
from pptx import Presentation
from pptx.util import Cm

from api.CollectionImage import CollectionImage
from api.Geometry import Geometry

#pylint: disable=too-few-public-methods
class Powerpoint():
    """ Summary
        -------
        High level interface to interact with the python-pptx module in
        the context of this specific app.

        Methods
        -------
        create_presentation(images_list)
            Creates a presentation with all the CollectionImages in a
            list
        to_slide(image)
            Creates a Powerpoint slide from a CollectionImage
    """

    @classmethod
    def create_presentation(cls, images_list, work_dir, output_path):
        """ Summary
            -------
            Creates a Powerpoint presentation with all the
            CollectionImages in a list

            Arguments
            ---------
            images_list : list(CollectionImage)
                A list of CollectionImages to export to pptx slides.
            work_dir : str
                Absolute path to the directory containing the images, so
                we can find the images on the user's disk.
            output_path : str
                Absolute path to where to put the generated Powerpoint
        """
        if not all(isinstance(i, CollectionImage) for i in images_list):
            raise TypeError("images_list must contain only CollectionImages")

        # create presentation
        prs = Presentation()

        for image in images_list:
            img_path = os.path.join(work_dir, image.filename)

            # create blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # compute layout parameters
            params = cls._get_layout_params(img_path)

            # create an image using these parameters

            left    =   params["image"]["left"]
            top     =   params["image"]["top"]
            height  =   params["image"]["height"]

            slide.shapes.add_picture(img_path, left, top, height=height)

            # create a legend text box using these parameters

            left    =   params["text"]["left"]
            top     =   params["text"]["top"]
            height  =   params["text"]["height"]
            width   =   params["text"]["width"]

            text_box = slide.shapes.add_textbox(left, top, width, height)

            # here the pptx interface is a bit ugly to use.
            # but basically our text frame contains all the text
            # separated in paragraphs
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            # this removes all paragraphs present except one empty
            # paragraph
            text_frame.clear()
            # so we can write in this one paragraph
            paragraph = text_frame.paragraphs[0]

            legend = image.to_legend()

            # to apply formatting such as italics defined between [i][/i]
            # tags, we must decompose the text and reassemble the segments
            # as pptx' run to be able to apply italics 
            # format italics
            # NOTE: supports only ONE set of italics tags
            italics = re.search(r"\[i\](.*)\[/i\]", legend)
            if italics:
                start_run = paragraph.add_run()
                start_run.text = legend[:italics.start()]
                
                italics_run = paragraph.add_run()
                italics_run.text = italics.group(1)
                font = italics_run.font
                font.italic = True

                end_run = paragraph.add_run()
                end_run.text = legend[italics.end():]
            
            else:
                text_frame.paragraphs[0].text = image.to_legend()

            # autofit text to the text box
            # text_frame.fit_text(
            #     font_family='Arial',
            #     max_size=18,
            #     bold=False,
            #     italic=False
            # )

        prs.save(output_path)


    @staticmethod
    def _get_layout_params(img_path, margin=1, textbox_height=1.5):
        """ Summary
            -------
            Calculates the top, left, and height parameters of our image
            and our text box to be given to the presentation.

            Arguments
            ---------
            img_path : str
                Absolute path to our image
            margin : float
                Margin size in cm
            textbox_height : float
                Height of the legend text box in cm

            Returns
            -------
            dict
                dict containing the parameters top, left, and height for
                the image and top, left, height and width for the text
                box.
                Values are in python-pptx Cm()

            NOTE
            ----
            All calculations below are done in centimeters !
            The actual size of the image in pixels is irrelevant:
            the only thing we need to know to determine the layout is
            its aspect ratio.
        """
        # dummy presentation to get access to slide width and height
        prs = Presentation()

        image_open = PIL.Image.open(img_path)
        image_size = image_open.size
        # canvas is the slide minus the margin and space for the text
        # values in cm
        canvas_width = prs.slide_width.cm - margin * 2
        canvas_height = prs.slide_height.cm - margin * 3 - textbox_height


        height, width = Geometry.fit_to_container(image_size, (canvas_width, canvas_height))

        top = margin + (canvas_height - height) / 2
        left = (prs.slide_width.cm - width) / 2

        params = {
            "image": {
                "left"  : Cm(left),
                "top"   : Cm(top),
                "height": Cm(height),
            },
            "text": {
                "left"  : Cm(margin),
                "top"   : Cm(margin * 2 + canvas_height),
                "width" : Cm(canvas_width),
                "height": Cm(textbox_height),
            }
        }

        return params


if __name__ == "__main__":
    ### Tests
    images = [
        CollectionImage(
            filename="Mona_Lisa,_by_Leonardo_da_Vinci,_from_C2RMF_retouched.JPG",
            artist="Sandro Botticelli",
            title="La Naissance de Vénus",
            datation="c. 1485",
            technique="Tempera sur toile",
            conservation_site="Gallerie Uffizi, Florence"
        ),
        CollectionImage(
            filename="Botticelli_Venus.jpg",
            artist="Sandro Botticelli",
            title="La Naissance de Vénus Mais avec un très long titre en espérant que ça va wrap",
            datation="c. 1485",
            technique="Tempera sur toile",
            conservation_site="Gallerie Uffizi, Florence"
        ),
        CollectionImage(
            filename="Strada_south_interior_facade.jpg",
            artist="Giulio Romano",
            title="Façade sud de la cour intérieure",
            datation="c. 1516",
            technique="Gravure",
            conservation_site="Palais du Te, Mantoue"
        )
    ]

    Powerpoint.create_presentation(
        images,
        "/Users/frieder/Documents/images",
        "test.pptx"
    )
